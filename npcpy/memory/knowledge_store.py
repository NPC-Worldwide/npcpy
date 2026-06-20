import hashlib
import os
import uuid
import yaml
import threading
from datetime import datetime, timezone
from typing import Dict, List, Any

from npcpy.llm_funcs import get_facts

DEFAULT_KNOWLEDGE_FILE = ".knowledge.yaml"


def _utcnow() -> str:
    return datetime.now(timezone.utc).isoformat()


def _make_id() -> str:
    return uuid.uuid4().hex


class KnowledgeStore:
    def __init__(self, directory: str, filename: str = None):
        self.directory = os.path.abspath(directory)
        self.filename = filename or DEFAULT_KNOWLEDGE_FILE
        self.file_path = os.path.join(self.directory, self.filename)
        self._lock = threading.RLock()

    def load(self) -> Dict[str, Any]:
        if not os.path.exists(self.file_path):
            return self._empty_template()
        with self._lock:
            try:
                with open(self.file_path, "r", encoding="utf-8") as f:
                    data = yaml.safe_load(f)
                if not isinstance(data, dict):
                    return self._empty_template()
                data.setdefault("memories", [])
                data.setdefault("knowledge", [])
                data.setdefault("concepts", [])
                data.setdefault("links", [])
                data.setdefault("scanned_files", {})
                return data
            except Exception:
                return self._empty_template()

    def save(self, data: Dict[str, Any]):
        data["updated_at"] = _utcnow()
        tmp = self.file_path + ".tmp"
        with self._lock:
            with open(tmp, "w", encoding="utf-8") as f:
                yaml.dump(data, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
            os.replace(tmp, self.file_path)

    def _save_no_index(self, data: Dict[str, Any]):
        data["updated_at"] = _utcnow()
        tmp = self.file_path + ".tmp"
        with self._lock:
            with open(tmp, "w", encoding="utf-8") as f:
                yaml.dump(data, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
            os.replace(tmp, self.file_path)

    def _empty_template(self) -> Dict[str, Any]:
        return {
            "created_at": _utcnow(),
            "updated_at": _utcnow(),
            "last_extracted_at": None,
            "last_evolved_at": None,
            "scanned_files": {},
            "memories": [],
            "knowledge": [],
            "concepts": [],
            "links": [],
        }

    def append_memory(self,
                      mem_id: str = None,
                      message_id: str = "",
                      conversation_id: str = "",
                      npc: str = "",
                      team: str = "",
                      directory_path: str = "",
                      initial_memory: str = "",
                      status: str = "pending_approval",
                      model: str = "",
                      provider: str = "",
                      final_memory: str = None,
                      source_type: str = "",
                      source_id: str = "",
                      **kwargs,
                      ) -> str:
        mem_id = mem_id or _make_id()
        mem = {
            "id": str(mem_id),
            "message_id": message_id,
            "conversation_id": conversation_id,
            "npc": npc,
            "team": team,
            "directory_path": directory_path,
            "timestamp": _utcnow(),
            "initial_memory": initial_memory,
            "final_memory": final_memory,
            "status": status,
            "model": model,
            "provider": provider,
            "created_at": _utcnow(),
            "source_type": source_type,
            "source_id": source_id,
        }
        mem.update(kwargs)
        data = self.load()
        data["memories"].append(mem)
        self.save(data)
        return mem_id

    def update_memory(self, mem_id: str, status: str, final_memory: str = None) -> bool:
        data = self.load()
        changed = False
        for mem in data.get("memories", []):
            if mem.get("id") == mem_id:
                mem["status"] = status
                if final_memory is not None:
                    mem["final_memory"] = final_memory
                changed = True
                break
        if changed:
            self.save(data)
        return changed

    def append_link(self, from_mem: str, to_mem: str, relation: str, agent: str = "") -> str:
        link = {
            "id": _make_id(),
            "from": from_mem,
            "to": to_mem,
            "relation": relation,
            "created_at": _utcnow(),
            "agent": agent,
        }
        data = self.load()
        data["knowledge"].append(link)
        self.save(data)
        return link["id"]

    def get_memories(self, status: str = None, limit: int = None) -> List[Dict[str, Any]]:
        data = self.load()
        mems = data.get("memories", [])
        if status:
            mems = [m for m in mems if m.get("status") == status]
        if limit:
            mems = mems[-limit:]
        return list(mems)

    def get_pending_memories(self) -> List[Dict[str, Any]]:
        return self.get_memories(status="pending_approval")

    def search_memories(self, query: str, limit: int = 20) -> List[Dict[str, Any]]:
        q = query.lower()
        results = []
        for mem in self.get_memories():
            text = (mem.get("initial_memory") or mem.get("final_memory") or "").lower()
            if q in text:
                results.append(mem)
                if limit and len(results) >= limit:
                    break
        return results

    def get_links(self) -> List[Dict[str, Any]]:
        return list(self.load().get("knowledge", []))

    def get_yaml_links(self) -> List[Dict[str, Any]]:
        return list(self.load().get("links", []))

    def get_concepts(self) -> List[Dict[str, Any]]:
        return list(self.load().get("concepts", []))

    def _extract_text(self, file_path: str) -> str:
        """Extract plain text from a supported document or code file."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                return "\n".join((p.extract_text() or "") for p in reader.pages)
            if ext == ".docx":
                from docx import Document
                doc = Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs)
            if ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                return "\n".join(parts)
            if ext in (".xlsx", ".xls"):
                import pandas as pd
                sheets = pd.read_excel(file_path, sheet_name=None)
                parts = []
                for name, df in sheets.items():
                    parts.append(f"Sheet: {name}\n{df.to_string(index=False)}")
                return "\n\n".join(parts)
            if ext == ".csv":
                import pandas as pd
                df = pd.read_csv(file_path)
                return df.to_string(index=False)
            if ext == ".html" or ext == ".htm":
                from bs4 import BeautifulSoup
                with open(file_path, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f, "html.parser")
                return soup.get_text(separator="\n", strip=True)
            if ext == ".json":
                import json
                with open(file_path, "r", encoding="utf-8") as f:
                    obj = json.load(f)
                return json.dumps(obj, indent=2, ensure_ascii=False)
            # Plain text / code fallback
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""

    @staticmethod
    def _is_noise_file(fpath: str, text: str) -> bool:
        """Heuristic: is this file just noise/boilerplate not worth a memory?"""
        name = os.path.basename(fpath).lower()
        # Lockfiles / manifests with no semantic value
        if name in {
            "package-lock.json", "yarn.lock", "pnpm-lock.yaml",
            "poetry.lock", "pipfile.lock", "cargo.lock", "gemfile.lock",
            "composer.lock", "go.sum", "requirements.lock",
        }:
            return True
        # Minified JS detection: average line length > 200 chars
        lines = text.splitlines()
        if len(lines) > 5:
            avg_len = sum(len(l) for l in lines) / len(lines)
            if avg_len > 200:
                return True
        # Too small after stripping — probably config or stub
        stripped = text.strip()
        if len(stripped) < 60:
            return True
        # Binary-looking content ratio (non-printable chars)
        printable = sum(1 for c in text if 32 <= ord(c) < 127 or c in "\n\r\t")
        if len(text) > 0 and printable / len(text) < 0.85:
            return True
        return False

    @staticmethod
    def _file_hash(fpath: str) -> str:
        """Fast SHA-256 hash of file contents."""
        h = hashlib.sha256()
        try:
            with open(fpath, "rb") as f:
                while True:
                    chunk = f.read(65536)
                    if not chunk:
                        break
                    h.update(chunk)
        except OSError:
            return ""
        return h.hexdigest()

    def extract_from_directory(
        self,
        include_extensions=None,
        exclude_dirs=None,
        max_chunk_size=2000,
        max_file_size_mb=5,
        model=None,
        provider=None,
        npc=None,
        context='',
    ) -> Dict[str, Any]:
        """Walk this store's directory, extract text from files, create memories.

        A ``.knowledge.yaml`` is created/updated in EACH directory that contains
        analyzed files.  Per-file state is tracked per-store so unchanged files are
        skipped.  Files whose content hasn't meaningfully changed are marked as
        noise and skipped until they change again.

        Facts are extracted via LLM on each file as it is processed.

        Prints progress to stdout and returns stats about what was extracted.
        """
        if include_extensions is None:
            include_extensions = {
                ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h",
                ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
                ".sh", ".bash", ".zsh", ".ps1", ".bat",
                ".txt", ".md", ".rst", ".log",
                ".csv", ".json", ".xml", ".yaml", ".yml",
                ".docx", ".pptx", ".xlsx", ".html", ".htm",
            }
        if exclude_dirs is None:
            exclude_dirs = {
                ".git", ".hg", ".svn", "node_modules", "__pycache__",
                ".pytest_cache", ".mypy_cache", "dist", "build",
                "venv", ".venv", "env", ".env", ".tox", ".egg-info",
                ".gitignore", ".gitattributes",
                ".local", ".cache", ".config", ".claude",
                ".github", ".vscode", ".idea", "vendor",
                "site-packages", "pip", "setuptools",
                "target", "out", ".gradle", ".terraform",
                "coverage", "htmlcov", ".nyc_output",
                "tmp", "temp", "logs", "uploads", "media",
                ".DS_Store", ".Trash", "__MACOSX",
                "third_party", "third-party", "3rdparty",
            }

        # Each subfolder that contains files gets its own .knowledge.yaml.
        # We load and save the store for EVERY FILE so a crash never loses progress.
        new_memories = 0
        files_scanned = 0
        files_skipped = 0
        files_noise = 0
        files_changed = 0
        stores_touched: set = set()

        max_bytes = max_file_size_mb * 1024 * 1024
        for root, dirs, files in os.walk(self.directory):
            # Prune excluded dirs in-place (also skip hidden dirs and site-packages anywhere)
            dirs[:] = [
                d for d in dirs
                if d not in exclude_dirs
                and not d.startswith(".")
                and "site-packages" not in d
                and "__pycache__" not in d
            ]

            for fname in files:
                # Skip hidden files, lock files, and the store itself
                if fname.startswith(".") or fname == self.filename:
                    continue
                ext = os.path.splitext(fname)[1].lower()
                if ext not in include_extensions:
                    continue

                fpath = os.path.join(root, fname)
                rel_path = os.path.relpath(fpath, self.directory)

                # Skip symlinks to avoid escaping the tree or looping
                if os.path.islink(fpath):
                    continue

                # Skip oversized files
                try:
                    fsize = os.path.getsize(fpath)
                except OSError:
                    continue
                if fsize > max_bytes:
                    continue

                try:
                    stat = os.stat(fpath)
                except OSError:
                    continue

                size = stat.st_size
                mtime = stat.st_mtime
                fhash = self._file_hash(fpath)

                # Load the store for the directory that contains this file
                sub_store = KnowledgeStore(root)
                sub_data = sub_store.load()
                sub_tracker = sub_data.setdefault("scanned_files", {})

                # Path inside the sub-store (usually just the filename)
                store_rel_path = os.path.relpath(fpath, root)

                prev = sub_tracker.get(store_rel_path, {})
                # Skip if unchanged
                if prev.get("hash") == fhash:
                    files_skipped += 1
                    print(f"  SKIP  {rel_path}")
                    continue

                files_changed += 1
                print(f"  SCAN  {rel_path}")
                text = self._extract_text(fpath)
                if not text or not text.strip():
                    sub_tracker[store_rel_path] = {
                        "hash": fhash,
                        "size": size,
                        "mtime": mtime,
                        "decision": "empty",
                        "scanned_at": _utcnow(),
                    }
                    sub_data["last_extracted_at"] = _utcnow()
                    sub_store.save(sub_data)
                    stores_touched.add(root)
                    print(f"  EMPTY {rel_path}")
                    continue

                if self._is_noise_file(fpath, text):
                    files_noise += 1
                    sub_tracker[store_rel_path] = {
                        "hash": fhash,
                        "size": size,
                        "mtime": mtime,
                        "decision": "noise",
                        "scanned_at": _utcnow(),
                    }
                    sub_data["last_extracted_at"] = _utcnow()
                    sub_store.save(sub_data)
                    stores_touched.add(root)
                    print(f"  NOISE {rel_path}")
                    continue

                files_scanned += 1
                print(f"  EXTRACT {rel_path}  ({size} bytes, {len(text)} chars)")

                # Extract facts via LLM — same chunking as kg_initial
                extracted_facts = []
                CHUNK = 100000
                if len(text) > CHUNK:
                    chunks = len(text) // CHUNK
                    print(f"    SLICING into {chunks} chunks (size {CHUNK})")
                    for n in range(chunks):
                        segment = text[n * CHUNK:(n + 1) * CHUNK]
                        print(f"    CHUNK {n + 1}/{chunks}  ({len(segment)} chars)")
                        chunk_facts = get_facts(segment, model=model, provider=provider, npc=npc, context=context)
                        if not chunk_facts:
                            print("    NO FACTS")
                        for fact in chunk_facts:
                            print(f"    FACT  {fact}")
                        extracted_facts.extend(chunk_facts)
                else:
                    extracted_facts = get_facts(text, model=model, provider=provider, npc=npc, context=context)
                    if not extracted_facts:
                        print("    NO FACTS")

                for fact in extracted_facts:
                    stmt = fact.get("statement", "").strip()
                    if not stmt:
                        continue
                    print(f"    MEM   {stmt[:200]}")
                    mem_id = _make_id()
                    mem = {
                        "id": mem_id,
                        "message_id": "",
                        "conversation_id": "",
                        "npc": "",
                        "team": "",
                        "directory_path": root,
                        "timestamp": _utcnow(),
                        "initial_memory": stmt,
                        "final_memory": stmt,
                        "status": "auto-extracted",
                        "model": "",
                        "provider": "",
                        "created_at": _utcnow(),
                        "source_type": "file",
                        "source_id": store_rel_path,
                        "source_hash": fhash,
                    }
                    sub_data["memories"].append(mem)
                    new_memories += 1

                sub_tracker[store_rel_path] = {
                    "hash": fhash,
                    "size": size,
                    "mtime": mtime,
                    "decision": "extracted",
                    "scanned_at": _utcnow(),
                    "facts": len(extracted_facts),
                }
                sub_data["last_extracted_at"] = _utcnow()
                sub_store.save(sub_data)
                stores_touched.add(root)

        return {
            "files_changed": files_changed,
            "files_skipped": files_skipped,
            "files_noise": files_noise,
            "files_scanned": files_scanned,
            "new_memories": new_memories,
            "stores_touched": sorted(stores_touched),
            "last_extracted_at": _utcnow(),
        }

    def get_links_for_memory(self, mem_id: str) -> Dict[str, List[Dict[str, Any]]]:
        data = self.load()
        out = []
        inn = []
        for link in data.get("knowledge", []):
            if link.get("from") == mem_id:
                out.append(link)
            if link.get("to") == mem_id:
                inn.append(link)
        return {"outgoing": out, "incoming": inn}

    def add_concept(self, name: str, description: str = "", memory_ids: List[str] = None) -> str:
        cid = _make_id()
        concept = {
            "id": cid,
            "name": name,
            "description": description,
            "memory_ids": list(memory_ids or []),
            "created_at": _utcnow(),
        }
        data = self.load()
        data["concepts"].append(concept)
        self.save(data)
        return cid

    def add_link(self, from_id: str, to_id: str, relation: str, link_type: str = "memory_to_memory") -> str:
        lid = _make_id()
        link = {
            "id": lid,
            "from": from_id,
            "to": to_id,
            "relation": relation,
            "type": link_type,
            "created_at": _utcnow(),
        }
        data = self.load()
        data["links"].append(link)
        self.save(data)
        return lid

    def build_context(self, max_memories: int = 10) -> str:
        parts = []
        mems = self.get_memories(status="human-approved", limit=max_memories)
        if mems:
            parts.append("Local knowledge:")
            for m in mems:
                parts.append(f"- {m.get('final_memory') or m.get('initial_memory')}")
        return "\n".join(parts) if parts else ""

    def evolve(self, model=None, provider=None, npc=None, context='',
               include_memories=True, include_knowledge=True, full_rebuild=False,
               all_facts=None, all_concepts=None) -> Dict[str, Any]:
        """Run concept extraction and linking on this store.

        If ``all_facts`` and ``all_concepts`` are supplied they are treated as the
        aggregated corpus (e.g. from multiple stores) so that cross-store linking
        works while the results are still persisted into *this* store.

        Extraction runs automatically first: only files whose content hash has
        changed since the last scan are re-evaluated.  Noise files (lockfiles,
        minified JS, tiny stubs) are skipped until they change again.
        """
        from npcpy.memory.knowledge_graph import kg_evolve_incremental

        extraction_stats = self.extract_from_directory(model=model, provider=provider, npc=npc, context=context)

        data = self.load()
        memories = data.get("memories", [])

        # 1. Build facts from memories (preserving memory IDs)
        facts = []
        stmt_to_mem_ids = {}
        for mem in memories:
            stmt = mem.get("final_memory") or mem.get("initial_memory", "")
            if stmt:
                mid = mem.get("id")
                facts.append({
                    "statement": stmt,
                    "source_text": stmt,
                    "type": mem.get("source_type", "memory"),
                    "generation": 0,
                    "memory_id": mid,
                })
                stmt_to_mem_ids.setdefault(stmt, []).append(mid)

        # Allow caller to inject an already-aggregated corpus
        if all_facts is not None:
            facts = list(all_facts)

        # 2. Build existing KG from current YAML state
        existing_concepts = all_concepts if all_concepts is not None else [
            {"name": c["name"], "description": c.get("description", ""), "generation": 0}
            for c in data.get("concepts", [])
        ]
        existing_kg = {
            "generation": 0,
            "facts": facts,
            "concepts": existing_concepts,
            "concept_links": [],
            "fact_to_concept_links": {},
            "fact_to_fact_links": [],
        }

        if not facts:
            return {"status": "skipped", "reason": "no_facts", "concepts_added": 0, "links_added": 0}

        # 3. Evolve in windows so we never blow the model context
        batch_size = 40
        new_kg = existing_kg
        for i in range(0, len(facts), batch_size):
            batch = facts[i:i + batch_size]
            print(f"  Evolving batch {i // batch_size + 1}/{(len(facts) - 1) // batch_size + 1} ({len(batch)} facts)...")
            new_kg, _ = kg_evolve_incremental(
                existing_kg=new_kg,
                new_facts=batch,
                model=model,
                provider=provider,
                npc=npc,
                context=context,
                get_concepts=True,
                link_concepts_facts=True,
                link_concepts_concepts=True,
                link_facts_facts=True,
            )

        # 4. Build concepts with stable IDs (reuse if name already exists)
        name_to_cid = {c["name"]: c["id"] for c in data.get("concepts", [])}
        yaml_concepts = data.get("concepts", []) if not full_rebuild else []
        for c in new_kg.get("concepts", []):
            cname = c["name"]
            if cname not in name_to_cid:
                name_to_cid[cname] = _make_id()
            cid = name_to_cid[cname]
            # Collect memory IDs linked to this concept
            mem_ids = set()
            for stmt, concept_names in new_kg.get("fact_to_concept_links", {}).items():
                if cname in concept_names:
                    for mid in stmt_to_mem_ids.get(stmt, []):
                        if mid:
                            mem_ids.add(mid)
            existing = next((x for x in yaml_concepts if x.get("name") == cname), None)
            if existing:
                existing["description"] = c.get("description", existing.get("description", ""))
                existing["generation"] = c.get("generation", existing.get("generation", 0))
                if mem_ids:
                    existing["memory_ids"] = list(set(existing.get("memory_ids", [])) | mem_ids)
            else:
                yaml_concepts.append({
                    "id": cid,
                    "name": cname,
                    "description": c.get("description", ""),
                    "generation": c.get("generation", 0),
                    "memory_ids": sorted(mem_ids),
                    "created_at": _utcnow(),
                })

        # 5. Build links: concept->memory, memory->memory, concept->concept
        seen_links = set()
        yaml_links = data.get("links", []) if not full_rebuild else []
        for link in yaml_links:
            seen_links.add((link.get("from"), link.get("to"), link.get("relation"), link.get("type")))

        # concept -> memory (fact_to_concept)
        for stmt, concept_names in new_kg.get("fact_to_concept_links", {}).items():
            for mid in stmt_to_mem_ids.get(stmt, []):
                if not mid:
                    continue
                for cname in concept_names:
                    cid = name_to_cid.get(cname)
                    if not cid:
                        continue
                    key = (mid, cid, "belongs_to", "memory_to_concept")
                    if key not in seen_links:
                        seen_links.add(key)
                        yaml_links.append({
                            "id": _make_id(),
                            "from": mid,
                            "to": cid,
                            "relation": "belongs_to",
                            "type": "memory_to_concept",
                            "created_at": _utcnow(),
                        })

        # memory -> memory (fact_to_fact)
        for s1, s2 in new_kg.get("fact_to_fact_links", []):
            for m1 in stmt_to_mem_ids.get(s1, []):
                if not m1:
                    continue
                for m2 in stmt_to_mem_ids.get(s2, []):
                    if not m2 or m1 == m2:
                        continue
                    key = tuple(sorted((m1, m2))) + ("related_to", "memory_to_memory")
                    if key not in seen_links:
                        seen_links.add(key)
                        yaml_links.append({
                            "id": _make_id(),
                            "from": m1,
                            "to": m2,
                            "relation": "related_to",
                            "type": "memory_to_memory",
                            "created_at": _utcnow(),
                        })

        for c1_name, c2_name in new_kg.get("concept_links", []):
            c1 = name_to_cid.get(c1_name)
            c2 = name_to_cid.get(c2_name)
            if not c1 or not c2 or c1 == c2:
                continue
            key = tuple(sorted((c1, c2))) + ("related_to", "concept_to_concept")
            if key not in seen_links:
                seen_links.add(key)
                yaml_links.append({
                    "id": _make_id(),
                    "from": c1,
                    "to": c2,
                    "relation": "related_to",
                    "type": "concept_to_concept",
                    "created_at": _utcnow(),
                })

        data["concepts"] = yaml_concepts
        data["links"] = yaml_links
        data["last_evolved_at"] = _utcnow()
        self.save(data)

        result = {
            "status": "success",
            "concepts_added": len(yaml_concepts),
            "links_added": len(yaml_links),
            "generation": new_kg.get("generation", 0),
            "last_evolved_at": data["last_evolved_at"],
        }
        if extraction_stats:
            result["extraction"] = extraction_stats
        return result

    def create(self, model=None, provider=None, npc=None, context="", content_text=None):
        """Build a fresh KG from store memories (and optional extra content), replacing existing concepts/links."""
        from npcpy.memory.knowledge_graph import kg_initial

        data = self.load()
        memories = data.get("memories", [])
        facts, stmt_to_mids = self._memories_to_facts(memories)

        if content_text:
            facts.append({"statement": content_text, "source_text": content_text, "type": "manual", "generation": 0, "memory_id": None})

        if not facts:
            return {"status": "skipped", "reason": "no_facts"}

        new_kg = kg_initial(
            content=None,
            facts=facts if facts else None,
            model=model,
            provider=provider,
            npc=npc,
            context=context,
            zoom_in_enabled=True,
        )

        stats = self._persist_kg_back(new_kg, stmt_to_mids, full_replace=True)
        return {"status": "success", "step": "create", **stats}

    def assimilate(self, model=None, provider=None, npc=None, context=""):
        """Alias for evolve: extract new files and merge into existing KG."""
        return self.evolve(model=model, provider=provider, npc=npc, context=context)

    def sleep(self, model=None, provider=None, npc=None, context="", operations=None):
        """Refine, prune, and deepen the existing KG."""
        from npcpy.memory.knowledge_graph import kg_sleep_process

        data = self.load()
        memories = data.get("memories", [])
        facts, stmt_to_mids = self._memories_to_facts(memories)
        existing_concepts = [
            {"name": c["name"], "description": c.get("description", ""), "generation": c.get("generation", 0)}
            for c in data.get("concepts", [])
        ]
        existing_kg = {
            "generation": 0,
            "facts": facts,
            "concepts": existing_concepts,
            "concept_links": [],
            "fact_to_concept_links": {},
            "fact_to_fact_links": [],
        }

        if not facts:
            return {"status": "skipped", "reason": "no_facts"}

        new_kg, _ = kg_sleep_process(
            existing_kg=existing_kg,
            model=model,
            provider=provider,
            npc=npc,
            context=context,
            operations_config=operations,
        )

        stats = self._persist_kg_back(new_kg, stmt_to_mids, full_replace=True)
        return {"status": "success", "step": "sleep", **stats}

    def dream(self, model=None, provider=None, npc=None, context="", num_seeds=3):
        """Creative synthesis from random concept seeds."""
        from npcpy.memory.knowledge_graph import kg_dream_process

        data = self.load()
        memories = data.get("memories", [])
        facts, stmt_to_mids = self._memories_to_facts(memories)
        existing_concepts = [
            {"name": c["name"], "description": c.get("description", ""), "generation": c.get("generation", 0)}
            for c in data.get("concepts", [])
        ]
        existing_kg = {
            "generation": 0,
            "facts": facts,
            "concepts": existing_concepts,
            "concept_links": [],
            "fact_to_concept_links": {},
            "fact_to_fact_links": [],
        }

        if len(existing_concepts) < num_seeds:
            return {"status": "skipped", "reason": "not_enough_concepts"}

        new_kg, _ = kg_dream_process(
            existing_kg=existing_kg,
            model=model,
            provider=provider,
            npc=npc,
            context=context,
            num_seeds=num_seeds,
        )

        stats = self._persist_kg_back(new_kg, stmt_to_mids, full_replace=False)
        return {"status": "success", "step": "dream", **stats}

    def _memories_to_facts(self, memories):
        facts = []
        stmt_to_mids = {}
        for mem in memories:
            stmt = (mem.get("final_memory") or mem.get("initial_memory", "")).strip()
            if not stmt:
                continue
            mid = mem.get("id")
            facts.append({
                "statement": stmt,
                "source_text": stmt,
                "type": mem.get("source_type", "memory"),
                "generation": 0,
                "memory_id": mid,
            })
            stmt_to_mids.setdefault(stmt, []).append(mid)
        return facts, stmt_to_mids

    def _persist_kg_back(self, new_kg, stmt_to_mids, full_replace=False):
        data = self.load()
        name_to_cid = {c["name"]: c["id"] for c in data.get("concepts", [])}
        yaml_concepts = [] if full_replace else list(data.get("concepts", []))

        for c in new_kg.get("concepts", []):
            cname = c["name"]
            if cname not in name_to_cid:
                name_to_cid[cname] = _make_id()
            cid = name_to_cid[cname]
            mem_ids = set()
            for stmt, concept_names in new_kg.get("fact_to_concept_links", {}).items():
                if cname in concept_names:
                    for mid in stmt_to_mids.get(stmt, []):
                        if mid:
                            mem_ids.add(mid)
            existing = next((x for x in yaml_concepts if x.get("name") == cname), None)
            if existing:
                existing["description"] = c.get("description", existing.get("description", ""))
                existing["generation"] = c.get("generation", existing.get("generation", 0))
                if mem_ids:
                    existing["memory_ids"] = list(set(existing.get("memory_ids", [])) | mem_ids)
            else:
                yaml_concepts.append({
                    "id": cid,
                    "name": cname,
                    "description": c.get("description", ""),
                    "generation": c.get("generation", 0),
                    "memory_ids": sorted(mem_ids),
                    "created_at": _utcnow(),
                })

        seen_links = set()
        yaml_links = [] if full_replace else list(data.get("links", []))
        for link in yaml_links:
            seen_links.add((link.get("from"), link.get("to"), link.get("relation"), link.get("type")))

        for stmt, concept_names in new_kg.get("fact_to_concept_links", {}).items():
            for mid in stmt_to_mids.get(stmt, []):
                if not mid:
                    continue
                for cname in concept_names:
                    cid = name_to_cid.get(cname)
                    if not cid:
                        continue
                    key = (mid, cid, "belongs_to", "memory_to_concept")
                    if key not in seen_links:
                        seen_links.add(key)
                        yaml_links.append({
                            "id": _make_id(),
                            "from": mid,
                            "to": cid,
                            "relation": "belongs_to",
                            "type": "memory_to_concept",
                            "created_at": _utcnow(),
                        })

        for s1, s2 in new_kg.get("fact_to_fact_links", []):
            for m1 in stmt_to_mids.get(s1, []):
                if not m1:
                    continue
                for m2 in stmt_to_mids.get(s2, []):
                    if not m2 or m1 == m2:
                        continue
                    key = tuple(sorted((m1, m2))) + ("related_to", "memory_to_memory")
                    if key not in seen_links:
                        seen_links.add(key)
                        yaml_links.append({
                            "id": _make_id(),
                            "from": m1,
                            "to": m2,
                            "relation": "related_to",
                            "type": "memory_to_memory",
                            "created_at": _utcnow(),
                        })

        for c1_name, c2_name in new_kg.get("concept_links", []):
            c1 = name_to_cid.get(c1_name)
            c2 = name_to_cid.get(c2_name)
            if not c1 or not c2 or c1 == c2:
                continue
            key = tuple(sorted((c1, c2))) + ("related_to", "concept_to_concept")
            if key not in seen_links:
                seen_links.add(key)
                yaml_links.append({
                    "id": _make_id(),
                    "from": c1,
                    "to": c2,
                    "relation": "related_to",
                    "type": "concept_to_concept",
                    "created_at": _utcnow(),
                })

        data["concepts"] = yaml_concepts
        data["links"] = yaml_links
        data["last_evolved_at"] = _utcnow()
        self.save(data)
        return {"concepts_added": len(yaml_concepts), "links_added": len(yaml_links)}

    @staticmethod
    def find_all(root_directory: str) -> List["KnowledgeStore"]:
        stores = []
        for dirpath, _dirnames, filenames in os.walk(root_directory):
            if DEFAULT_KNOWLEDGE_FILE in filenames:
                stores.append(KnowledgeStore(dirpath))
        return stores

    @staticmethod
    def aggregate(root_directory: str, max_depth: int = None) -> Dict[str, Any]:
        stores = KnowledgeStore.find_all(root_directory)
        if max_depth is not None:
            filtered = []
            for s in stores:
                rel = os.path.relpath(s.directory, root_directory).split(os.sep)
                if len(rel) <= max_depth:
                    filtered.append(s)
            stores = filtered
        all_memories = []
        all_knowledge = []
        for s in stores:
            d = s.load()
            all_memories.extend(d.get("memories", []))
            all_knowledge.extend(d.get("knowledge", []))
        return {
            "root": root_directory,
            "sources": [s.directory for s in stores],
            "memories": all_memories,
            "knowledge": all_knowledge,
        }

    @staticmethod
    def init_directory(directory: str) -> "KnowledgeStore":
        store = KnowledgeStore(directory)
        if not os.path.exists(store.file_path):
            store.save(store._empty_template())
        return store


def get_store_for_path(path: str) -> KnowledgeStore:
    return KnowledgeStore.init_directory(path)
